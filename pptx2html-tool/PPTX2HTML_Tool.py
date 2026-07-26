#!/usr/bin/env python3
"""
pptx2html v2 — transforme un .pptx (Google Slides / PowerPoint) en une page
HTML unique : visionneuse + ÉDITEUR INTERACTIF intégré.

Le HTML généré s'ouvre comme un site : bouton ✏️ pour passer en mode édition
(dessiner des boutons/zones cliquables, cacher des diapos pour qu'elles ne
soient accessibles que via des boutons, incruster des vidéos mp4 ou YouTube),
puis 💾 Enregistrer re-télécharge le fichier mis à jour — aucun serveur,
100 % local, rien ne transite par un service tiers.

Usage : python PPTX2HTML_Tool.py deck.pptx [-o sortie.html] [--dpi 150] [--assets]

Pipeline : pptx -> PDF (LibreOffice) -> images JPEG (PyMuPDF, repli Poppler)
-> HTML unique. Les liens posés dans Slides (sauts de diapo, URLs, YouTube)
sont importés automatiquement comme zones interactives, les vidéos mp4
embarquées sont extraites et posées sur les diapos.
"""
import argparse, base64, glob, html, json, os, re, shutil, subprocess, sys, tempfile

APP_VERSION = "2.0.0"
SOFFICE_WRAPPER = "/mnt/skills/public/pptx/scripts/office/soffice.py"


def app_dir():
    """Dossier de l'application (exe PyInstaller ou script)."""
    if getattr(sys, "frozen", False):
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))


class ConvError(Exception):
    pass


def find_soffice():
    """Trouve LibreOffice selon l'OS."""
    if os.path.isfile(SOFFICE_WRAPPER):
        return [sys.executable, SOFFICE_WRAPPER]
    cands = [shutil.which("soffice"),
             r"C:\Program Files\LibreOffice\program\soffice.exe",
             r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
             "/Applications/LibreOffice.app/Contents/MacOS/soffice"]
    for c in cands:
        if c and os.path.isfile(c):
            return [c]
    raise ConvError("LibreOffice introuvable. Installe-le depuis libreoffice.org\n"
                    "(ou ajoute soffice au PATH).")


def get_pymupdf():
    """PyMuPDF si disponible (rendu PDF→images sans binaire externe —
    compatible Smart App Control). Retourne le module ou None."""
    try:
        import pymupdf
        return pymupdf
    except ImportError:
        pass
    try:
        import fitz  # ancien nom du module PyMuPDF
        return fitz
    except ImportError:
        return None


def find_pdftoppm():
    # Poppler livré à côté de l'exe (zip Windows) — prioritaire
    exe = "pdftoppm.exe" if sys.platform == "win32" else "pdftoppm"
    bundled = os.path.join(app_dir(), "poppler", exe)
    if os.path.isfile(bundled):
        return bundled
    c = shutil.which("pdftoppm")
    if c:
        return c
    for extra in [r"C:\poppler\Library\bin\pdftoppm.exe",
                  r"C:\Program Files\poppler\Library\bin\pdftoppm.exe",
                  "/opt/homebrew/bin/pdftoppm", "/usr/local/bin/pdftoppm"]:
        if os.path.isfile(extra):
            return extra
    raise ConvError("Poppler (pdftoppm) introuvable.\n"
                    "Windows : télécharger poppler et le dézipper dans C:\\poppler\n"
                    "Mac : brew install poppler")


def run(cmd, **kw):
    if sys.platform == "win32":
        # pas de fenêtre console qui clignote depuis l'exe/pythonw
        kw.setdefault("creationflags", subprocess.CREATE_NO_WINDOW)
    r = subprocess.run(cmd, capture_output=True, text=True, **kw)
    if r.returncode != 0:
        raise ConvError(f"Échec : {' '.join(map(str, cmd))}\n{r.stderr[:500]}")
    return r


# ========================= EXTRACTION PPTX =========================

YT_RE = re.compile(
    r"(?:youtube\.com/(?:watch\?\S*?v=|shorts/|embed/|live/)|youtu\.be/)([\w-]{6,})")


def zone_from_url(url):
    """Un lien YouTube (export Google Slides d'une vidéo) devient une zone
    « lire la vidéo » en overlay ; tout autre lien reste une zone URL."""
    m = YT_RE.search(url or "")
    if m:
        return {"action": "video",
                "video": {"url": "https://www.youtube.com/embed/" + m.group(1)}}
    return {"action": "url", "url": url}


def extract_deck(pptx_path, log=print):
    """Une passe python-pptx : notes, zones cliquables (les liens posés dans
    Slides/PowerPoint deviennent des actions) et vidéos embarquées.
    Retourne (notes, zones, videos, media)."""
    try:
        from pptx import Presentation
        from pptx.enum.action import PP_ACTION
        from pptx.oxml.ns import qn
    except ImportError:
        log("python-pptx absent : notes, liens et vidéos non extraits.")
        return [], [], [], {}
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        log(f"Lecture du pptx impossible ({e!r}) : conversion images seule.")
        return [], [], [], {}

    sw, sh = prs.slide_width, prs.slide_height
    slide_index = {s.slide_id: i for i, s in enumerate(prs.slides)}
    n_slides = len(slide_index)
    playable = {"video/mp4", "video/webm", "video/ogg", "video/quicktime"}
    notes, zones, videos, media, seen = [], [], [], {}, {}

    for slide in prs.slides:
        txt = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf is not None:
                txt = tf.text.strip()
        notes.append(txt)

        zs, vs = [], []
        for shape in slide.shapes:
            if None in (shape.left, shape.top, shape.width, shape.height):
                continue
            rect = {"x": round(shape.left / sw * 100, 2),
                    "y": round(shape.top / sh * 100, 2),
                    "w": round(shape.width / sw * 100, 2),
                    "h": round(shape.height / sh * 100, 2)}

            # --- vidéo embarquée dans le pptx ---
            vf = shape._element.findall(".//" + qn("a:videoFile"))
            if vf:
                rid = vf[0].get(qn("r:link"))
                rel = slide.part.rels.get(rid) if rid else None
                if rel is not None and not rel.is_external:
                    part = rel.target_part
                    key = str(part.partname)
                    if key in seen:
                        mid = seen[key]
                    else:
                        ct = (part.content_type or "").lower()
                        if ct not in playable:
                            log(f"Vidéo ignorée (format non lisible en HTML : {ct})")
                            mid = None
                        else:
                            mid = "v%d" % (len(media) + 1)
                            media[mid] = {
                                "mime": "video/mp4" if ct == "video/quicktime" else ct,
                                "data": base64.b64encode(part.blob).decode("ascii")}
                        seen[key] = mid
                    if mid:
                        vs.append(dict(rect, media=mid, controls=True))
                continue  # une vidéo n'est pas aussi une zone cliquable

            # --- lien posé sur la forme entière ---
            act = None
            try:
                ca = shape.click_action
                if ca.target_slide is not None:
                    act = {"action": "goto",
                           "slide": slide_index[ca.target_slide.slide_id]}
                elif ca.action == PP_ACTION.NEXT_SLIDE:
                    act = {"action": "next"}
                elif ca.action == PP_ACTION.PREVIOUS_SLIDE:
                    act = {"action": "prev"}
                elif ca.action == PP_ACTION.FIRST_SLIDE:
                    act = {"action": "goto", "slide": 0}
                elif ca.action == PP_ACTION.LAST_SLIDE:
                    act = {"action": "goto", "slide": n_slides - 1}
                elif ca.hyperlink.address:
                    act = zone_from_url(ca.hyperlink.address)
            except Exception:
                pass
            # --- sinon, lien sur un run de texte dans la forme ---
            if act is None and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run_ in para.runs:
                        try:
                            if run_.hyperlink.address:
                                act = zone_from_url(run_.hyperlink.address)
                                break
                        except Exception:
                            pass
                    if act:
                        break
            if act:
                zs.append(dict(rect, **act))
        zones.append(zs)
        videos.append(vs)
    return notes, zones, videos, media


# ========================= CONVERSION =========================

def convert(pptx_path, out_html, dpi=150, lang="fr", embed=True, log=print):
    pptx_path = os.path.abspath(pptx_path)
    if not os.path.isfile(pptx_path):
        raise ConvError(f"Fichier introuvable : {pptx_path}")

    title = os.path.splitext(os.path.basename(pptx_path))[0]
    tmp = tempfile.mkdtemp(prefix="pptx2html_")
    try:
        # 1) pptx -> pdf
        log("Rendu des slides via LibreOffice…")
        run(find_soffice() + ["--headless", "--convert-to", "pdf",
                              "--outdir", tmp, pptx_path])
        pdfs = glob.glob(os.path.join(tmp, "*.pdf"))
        if not pdfs:
            raise ConvError("La conversion PDF a échoué (aucun PDF produit).")
        pdf = pdfs[0]

        # 2) pdf -> jpegs (PyMuPDF en priorité, Poppler en repli)
        log("Génération des images…")
        mu = get_pymupdf()
        if mu is not None:
            doc = mu.open(pdf)
            imgs = []
            for i, page in enumerate(doc, 1):
                p = os.path.join(tmp, f"slide-{i:03d}.jpg")
                page.get_pixmap(dpi=dpi).save(p, jpg_quality=85)
                imgs.append(p)
            doc.close()
        else:
            run([find_pdftoppm(), "-jpeg", "-jpegopt", "quality=85",
                 "-r", str(dpi), pdf, os.path.join(tmp, "slide")])
            imgs = sorted(glob.glob(os.path.join(tmp, "slide-*.jpg")))
        if not imgs:
            raise ConvError("Aucune image générée.")

        # 3) base64 (embarqué) ou dossier d'assets (non portable)
        slides_src, assets, assets_dir = [], None, None
        if embed:
            for p in imgs:
                with open(p, "rb") as f:
                    slides_src.append(base64.b64encode(f.read()).decode("ascii"))
        else:
            assets = os.path.splitext(os.path.basename(out_html))[0] + "_assets"
            assets_dir = os.path.join(os.path.dirname(os.path.abspath(out_html)), assets)
            os.makedirs(assets_dir, exist_ok=True)
            for i, p in enumerate(imgs, 1):
                name = f"slide-{i:03d}.jpg"
                shutil.copy(p, os.path.join(assets_dir, name))
                slides_src.append(f"{assets}/{name}")

        # 4) extraction pptx : notes, liens -> zones, vidéos embarquées
        notes, zones, videos, media = extract_deck(pptx_path, log=log)
        n = len(slides_src)
        notes = (notes + [""] * n)[:n]
        zones = (zones + [[] for _ in range(n)])[:n]
        videos = (videos + [[] for _ in range(n)])[:n]
        nz = sum(len(z) for z in zones)
        nv = sum(len(v) for v in videos)
        if nz:
            log(f"{nz} zone(s) cliquable(s) importée(s) depuis le pptx")
        if nv:
            log(f"{nv} vidéo(s) embarquée(s) extraite(s)")

        if not embed and media:
            # mode assets : vidéos posées en fichiers à côté des images
            exts = {"video/mp4": ".mp4", "video/webm": ".webm", "video/ogg": ".ogv"}
            for mid, m in media.items():
                name = "media-%s%s" % (mid, exts.get(m["mime"], ".bin"))
                with open(os.path.join(assets_dir, name), "wb") as fh:
                    fh.write(base64.b64decode(m.pop("data")))
                m["path"] = "%s/%s" % (assets, name)

        # 5) HTML unique (visionneuse + éditeur)
        write_html(out_html, title, slides_src, notes, zones, videos, media,
                   lang, embed)
        size_mb = os.path.getsize(out_html) / 1e6
        mode = "tout embarqué" if embed else "assets séparés (non portable)"
        log(f"Terminé : {len(slides_src)} slides, {size_mb:.1f} Mo ({mode})")
        log("Ouvre le HTML et clique ✏️ (ou touche E) : boutons, diapos "
            "cachées, vidéos — puis 💾 Enregistrer.")
        return out_html
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


# ========================= HTML (visionneuse + éditeur) =========================
# Le fichier généré ne contient que 3 balises <script> : la config (JSON),
# les assets (JSON) et l'application (ce JS). L'application reconstruit le
# document entier à partir de ces 3 balises pour s'enregistrer elle-même
# (bouton 💾) — aucun serveur, tout reste local.

APP_JS = r"""
(function () {
'use strict';
var $ = function (id) { return document.getElementById(id); };
var CFG = JSON.parse($('cfg').textContent);
var ASSETS = JSON.parse($('assets').textContent);
var META = CFG.meta, SLIDES = CFG.slides;
SLIDES.forEach(function (s) { s.zones = s.zones || []; s.videos = s.videos || []; s.notes = s.notes || ''; });
var FR = META.lang !== 'en';
var IMG = function (i) { return META.embed ? 'data:image/jpeg;base64,' + ASSETS.images[i] : ASSETS.images[i]; };
var MEDIA = function (id) { var m = ASSETS.media[id]; return m ? (m.data ? 'data:' + m.mime + ';base64,' + m.data : m.path) : ''; };
var T = FR ? { of: '/', notes: 'Notes', noNotes: 'Aucune note pour cette diapo.',
  back: '\u21a9 Retour', hiddenBadge: 'Diapo cach\u00e9e', slide: 'diapo',
  help: '\u2190 \u2192 : naviguer \u00b7 F : plein \u00e9cran \u00b7 N : notes \u00b7 T : vignettes' }
: { of: '/', notes: 'Notes', noNotes: 'No notes for this slide.',
  back: '\u21a9 Back', hiddenBadge: 'Hidden slide', slide: 'slide',
  help: '\u2190 \u2192 : navigate \u00b7 F : fullscreen \u00b7 N : notes \u00b7 T : thumbnails' };

document.title = META.title;

/* ---------- styles ---------- */
var style = document.createElement('style');
style.textContent = "\
:root{--bg:#111318;--panel:#1b1e26;--panel2:#232733;--fg:#e8eaf0;--muted:#8b90a0;--accent:#5b8cff;--warn:#ffb020;--radius:10px}\
*{margin:0;padding:0;box-sizing:border-box}\
html,body{height:100%;background:var(--bg);color:var(--fg);font:15px/1.5 -apple-system,'Segoe UI',Roboto,Helvetica,Arial,sans-serif;overflow:hidden}\
.hidden{display:none!important}\
#app{display:flex;flex-direction:column;height:100%}\
header{display:flex;align-items:center;gap:10px;padding:10px 16px;background:var(--panel);border-bottom:1px solid #2a2e3a;flex-shrink:0}\
header h1{font-size:15px;font-weight:600;white-space:nowrap;overflow:hidden;text-overflow:ellipsis;flex:1}\
body.editing header h1{cursor:text}\
#counter{color:var(--muted);font-variant-numeric:tabular-nums;white-space:nowrap}\
button.icon{background:var(--panel2);color:var(--fg);border:none;border-radius:8px;padding:7px 12px;cursor:pointer;font-size:14px;transition:background .15s;white-space:nowrap}\
button.icon:hover{background:#2e3342}\
button.icon.active{background:var(--accent);color:#fff}\
#tools{display:flex;gap:8px}\
#main{flex:1;display:flex;min-height:0}\
#stage{flex:1;display:flex;align-items:center;justify-content:center;position:relative;padding:18px;min-width:0}\
#wrap{position:relative;max-width:100%;max-height:100%;display:flex;transition:opacity .18s ease}\
#wrap.fading{opacity:0}\
#slide{max-width:100%;max-height:100%;border-radius:var(--radius);box-shadow:0 8px 40px rgba(0,0,0,.55);user-select:none;display:block}\
body.editing{user-select:none}\
body.drawing #stage{cursor:crosshair}\
.navzone{position:absolute;top:0;bottom:0;width:22%;cursor:pointer;display:flex;align-items:center;opacity:0;transition:opacity .2s;z-index:1}\
.navzone:hover{opacity:1}\
.navzone span{font-size:34px;color:#fff;background:rgba(0,0,0,.45);border-radius:50%;width:52px;height:52px;display:flex;align-items:center;justify-content:center}\
#prev{left:0;justify-content:flex-start;padding-left:14px}\
#next{right:0;justify-content:flex-end;padding-right:14px}\
body.editing .navzone,#stage.onhidden .navzone{display:none}\
.zone{position:absolute;z-index:3;border-radius:6px}\
.look-hover{cursor:pointer;border:2px solid transparent;transition:border-color .15s,background .15s}\
.look-hover:hover{border-color:var(--accent);background:rgba(91,140,255,.12)}\
.look-outline{cursor:pointer;border:2px solid var(--accent)}\
.look-button{cursor:pointer;display:flex;align-items:center;justify-content:center;background:var(--accent);color:#fff;font-weight:600;border-radius:10px;box-shadow:0 3px 14px rgba(0,0,0,.4);text-align:center;overflow:hidden;padding:2px 8px}\
.look-button:hover{filter:brightness(1.12)}\
body.editing .zone{outline:1px dashed rgba(91,140,255,.9)}\
body.editing .ov{cursor:grab}\
body.editing .ov.sel{outline:2px solid var(--warn)}\
.hdl{position:absolute;right:-7px;bottom:-7px;width:14px;height:14px;background:var(--warn);border:2px solid #111;border-radius:50%;cursor:nwse-resize;z-index:6}\
.vidbox{position:absolute;z-index:2;background:#000;border-radius:8px}\
.vidbox video,.vidbox iframe{width:100%;height:100%;border:0;border-radius:8px;display:block}\
.vcover{position:absolute;inset:0;z-index:4}\
#backBtn{position:absolute;top:16px;left:16px;z-index:6;background:rgba(20,22,30,.85);color:#fff;border:1px solid #2a2e3a;border-radius:20px;padding:8px 16px;cursor:pointer;font-size:14px}\
#backBtn:hover{background:var(--accent);border-color:var(--accent)}\
#hidBadge{position:absolute;top:16px;right:16px;z-index:6;background:rgba(192,57,43,.85);color:#fff;border-radius:14px;padding:4px 12px;font-size:12px}\
#props{width:280px;background:var(--panel);border-left:1px solid #2a2e3a;padding:14px;overflow-y:auto;flex-shrink:0;font-size:13px}\
#props h3{font-size:12px;text-transform:uppercase;letter-spacing:.6px;color:var(--muted);margin:4px 0 10px}\
#props label{display:block;margin:10px 0 2px}\
#props label.ck{display:flex;gap:8px;align-items:flex-start;cursor:pointer}\
#props input[type=text],#props select{width:100%;background:var(--panel2);color:var(--fg);border:1px solid #2a2e3a;border-radius:6px;padding:6px 8px;font-size:13px;margin-top:4px}\
#props input[type=color]{margin-top:4px;width:52px;height:30px;border:none;background:none;cursor:pointer}\
#props hr{border:none;border-top:1px solid #2a2e3a;margin:14px 0}\
#props .muted{color:var(--muted);line-height:1.55}\
#props button.danger{margin-top:14px;background:#3a2326;color:#ff8a8a;border:none;border-radius:8px;padding:8px 12px;cursor:pointer;width:100%;font-size:13px}\
#props button.danger:hover{background:#54282d}\
#thumbs{width:168px;background:var(--panel);border-left:1px solid #2a2e3a;overflow-y:auto;padding:10px;display:flex;flex-direction:column;gap:8px;flex-shrink:0}\
.th{position:relative;cursor:pointer}\
.th img{width:100%;display:block;border-radius:6px;border:2px solid transparent;opacity:.65;transition:.15s}\
.th:hover img{opacity:1}\
.th.current img{border-color:var(--accent);opacity:1}\
.th.th-hidden img{opacity:.28;filter:grayscale(.8)}\
.tnum{position:absolute;left:6px;bottom:6px;font-size:10px;background:rgba(0,0,0,.6);color:#fff;padding:1px 6px;border-radius:8px;pointer-events:none}\
.teye{position:absolute;top:4px;right:4px;background:rgba(0,0,0,.6);border:none;border-radius:6px;padding:2px 5px;cursor:pointer;font-size:12px}\
#notes{background:var(--panel);border-top:1px solid #2a2e3a;padding:12px 18px;max-height:26vh;overflow-y:auto;flex-shrink:0;white-space:pre-wrap;color:var(--muted);font-size:14px}\
#notes b{color:var(--fg);display:block;margin-bottom:4px}\
#progress{height:3px;background:var(--accent);width:0;transition:width .25s ease;flex-shrink:0}\
#hint{position:fixed;bottom:14px;left:50%;transform:translateX(-50%);background:rgba(20,22,30,.92);padding:8px 16px;border-radius:20px;font-size:12.5px;color:var(--muted);pointer-events:none;transition:opacity .5s;white-space:nowrap;z-index:20}\
#lightbox{position:fixed;inset:0;background:rgba(5,6,10,.9);z-index:50;display:flex;align-items:center;justify-content:center}\
#lb{width:min(92vw,1200px);aspect-ratio:16/9;background:#000;border-radius:10px;box-shadow:0 10px 60px rgba(0,0,0,.7)}\
#lb video,#lb iframe{width:100%;height:100%;border:0;border-radius:10px}\
#lbClose{position:fixed;top:18px;right:22px;z-index:51;background:rgba(255,255,255,.14);color:#fff;border:none;border-radius:50%;width:40px;height:40px;font-size:17px;cursor:pointer}\
#vmenu{position:fixed;z-index:40;background:var(--panel2);border:1px solid #2a2e3a;border-radius:10px;padding:6px;display:flex;flex-direction:column;gap:2px;box-shadow:0 8px 30px rgba(0,0,0,.5)}\
#vmenu button{background:none;border:none;color:var(--fg);padding:8px 12px;text-align:left;border-radius:6px;cursor:pointer;font-size:13px}\
#vmenu button:hover{background:var(--panel)}\
@media (max-width:700px){#thumbs{display:none}.navzone span{display:none}}";
document.head.appendChild(style);

/* ---------- squelette (ajouté APRES les balises script, jamais à leur place :
   elles servent à régénérer le fichier au moment d'enregistrer) ---------- */
document.body.insertAdjacentHTML('beforeend',
'<div id="app">' +
'<header>' +
'<h1 id="title"></h1>' +
'<span id="counter"></span>' +
'<span id="tools" class="hidden">' +
'<button class="icon" id="btnAddZone" title="Dessiner une zone cliquable sur la diapo">\u2795 Zone</button>' +
'<button class="icon" id="btnAddVideo" title="Incruster une vid\u00e9o sur la diapo">\ud83c\udfac Vid\u00e9o</button>' +
'<button class="icon" id="btnSave" title="T\u00e9l\u00e9charger ce fichier mis \u00e0 jour">\ud83d\udcbe Enregistrer</button>' +
'<button class="icon" id="btnLock" title="Version verrouill\u00e9e (sans mode \u00e9dition) pour diffusion">\ud83d\udd12 Export final</button>' +
'</span>' +
(META.locked ? '' : '<button class="icon" id="btnEdit" title="E">\u270f\ufe0f</button>') +
'<button class="icon hidden" id="btnNotes" title="N">\ud83d\uddd2 ' + T.notes + '</button>' +
'<button class="icon" id="btnThumbs" title="T">\u25a6</button>' +
'<button class="icon" id="btnFS" title="F">\u26f6</button>' +
'</header>' +
'<div id="main">' +
'<div id="stage">' +
'<div id="wrap"><img id="slide" alt=""></div>' +
'<div class="navzone" id="prev"><span>\u2039</span></div>' +
'<div class="navzone" id="next"><span>\u203a</span></div>' +
'<button id="backBtn" class="hidden"></button>' +
'<span id="hidBadge" class="hidden"></span>' +
'</div>' +
'<aside id="props" class="hidden"></aside>' +
'<div id="thumbs"></div>' +
'</div>' +
'<div id="notes" class="hidden"></div>' +
'<div id="progress"></div>' +
'</div>' +
'<div id="lightbox" class="hidden"><div id="lb"></div></div>' +
'<button id="lbClose" class="hidden">\u2715</button>' +
'<div id="vmenu" class="hidden">' +
'<button id="vmFile">\ud83d\udcc1 Fichier vid\u00e9o local (lecture hors ligne)</button>' +
'<button id="vmYt">\u25b6 Lien YouTube (n\u00e9cessite internet)</button>' +
'</div>' +
'<input type="file" id="filePick" accept="video/*" class="hidden">' +
'<div id="hint">' + T.help + (META.locked ? '' : ' \u00b7 E : \u00e9dition') + '</div>');

var wrap = $('wrap'), slideEl = $('slide'), counter = $('counter'),
    thumbs = $('thumbs'), notesEl = $('notes'), prog = $('progress'),
    props = $('props'), stage = $('stage'), backBtn = $('backBtn'),
    hidBadge = $('hidBadge'), titleEl = $('title'), btnNotes = $('btnNotes');
var cur = 0, editMode = false, drawMode = false, dirty = false,
    sel = null, drag = null, hist = [], thumbItems = [];
titleEl.textContent = META.title;
hidBadge.textContent = T.hiddenBadge;
backBtn.textContent = T.back;
slideEl.draggable = false;

/* ---------- helpers ---------- */
function esc(s) { return String(s == null ? '' : s).replace(/[&<>]/g, function (c) { return { '&': '&amp;', '<': '&lt;', '>': '&gt;' }[c]; }); }
function escA(s) { return esc(s).replace(/"/g, '&quot;'); }
function clamp(v, a, b) { return Math.max(a, Math.min(b, v)); }
function markDirty() { dirty = true; titleEl.textContent = META.title + ' \u2022'; }
function clearDirty() { dirty = false; titleEl.textContent = META.title; }
function visCount() { return SLIDES.filter(function (s) { return !s.hidden; }).length || 1; }
function visPos(i) { var p = 0; for (var k = 0; k <= i; k++) if (!SLIDES[k].hidden) p++; return p; }
function firstVisible() { for (var i = 0; i < SLIDES.length; i++) if (!SLIDES[i].hidden) return i; return 0; }
function lastVisible() { for (var i = SLIDES.length - 1; i >= 0; i--) if (!SLIDES[i].hidden) return i; return SLIDES.length - 1; }
function linNext() { for (var i = cur + 1; i < SLIDES.length; i++) if (editMode || !SLIDES[i].hidden) return i; return cur; }
function linPrev() { for (var i = cur - 1; i >= 0; i--) if (editMode || !SLIDES[i].hidden) return i; return cur; }
function relPct(e) {
  var r = wrap.getBoundingClientRect();
  return { x: clamp((e.clientX - r.left) / r.width * 100, 0, 100),
           y: clamp((e.clientY - r.top) / r.height * 100, 0, 100) };
}
function setRect(el, o) { el.style.left = o.x + '%'; el.style.top = o.y + '%'; el.style.width = o.w + '%'; el.style.height = o.h + '%'; }
function ytEmbed(u) {
  var m = String(u || '').match(/(?:youtube\.com\/(?:watch\?\S*?v=|shorts\/|embed\/|live\/)|youtu\.be\/)([\w-]{6,})/);
  return m ? 'https://www.youtube.com/embed/' + m[1] : u;
}

/* ---------- navigation ---------- */
function go(i, opts) {
  opts = opts || {};
  i = clamp(i, 0, SLIDES.length - 1);
  if (!opts.noHist && i !== cur) { hist.push(cur); if (hist.length > 200) hist.shift(); }
  cur = i;
  var apply = function () { slideEl.src = IMG(SLIDES[cur].img); renderOverlays(); wrap.classList.remove('fading'); };
  if (opts.instant) apply();
  else { wrap.classList.add('fading'); setTimeout(apply, 120); }
  var s = SLIDES[cur];
  if (editMode) counter.textContent = T.slide + ' ' + (cur + 1) + ' / ' + SLIDES.length;
  else if (s.hidden) counter.textContent = '\u2022';
  else {
    counter.textContent = visPos(cur) + ' ' + T.of + ' ' + visCount();
    prog.style.width = (visPos(cur) / visCount() * 100) + '%';
  }
  stage.classList.toggle('onhidden', !!s.hidden && !editMode);
  backBtn.classList.toggle('hidden', editMode || !s.hidden);
  hidBadge.classList.toggle('hidden', !(editMode && s.hidden));
  syncThumbs();
  notesEl.innerHTML = '<b>' + T.notes + ' \u2014 ' + T.slide + ' ' + (cur + 1) + '</b>' +
    (s.notes ? esc(s.notes) : '<i>' + T.noNotes + '</i>');
  location.hash = cur + 1;
  if (editMode) { if (!opts.keepSel) sel = null; renderProps(); }
}
function refresh() { go(cur, { instant: true, noHist: true, keepSel: true }); }
function goBack() { var p = hist.pop(); go(p == null ? firstVisible() : p, { noHist: true }); }

/* ---------- rendu des zones et vidéos ---------- */
function renderOverlays() {
  wrap.querySelectorAll('.ov').forEach(function (n) { n.remove(); });
  var s = SLIDES[cur];
  s.videos.forEach(function (v, vi) {
    var box = document.createElement('div');
    box.className = 'ov vidbox';
    setRect(box, v);
    var el;
    if (v.url) {
      el = document.createElement('iframe');
      el.src = v.url;
      el.allow = 'autoplay; fullscreen; encrypted-media; picture-in-picture';
      el.allowFullscreen = true;
    } else {
      el = document.createElement('video');
      el.src = MEDIA(v.media);
      el.controls = v.controls !== false;
      el.loop = !!v.loop;
      el.muted = !!v.muted || !!v.autoplay;
      el.autoplay = !!v.autoplay && !editMode;
      el.playsInline = true;
    }
    box.appendChild(el);
    if (editMode) {
      var cov = document.createElement('div');
      cov.className = 'vcover';
      box.appendChild(cov);
      if (sel && sel.kind === 'video' && sel.i === vi) { box.classList.add('sel'); addHandle(box); }
      attachEdit(box, 'video', vi, v);
    }
    wrap.appendChild(box);
  });
  s.zones.forEach(function (z, zi) {
    var d = document.createElement('div');
    d.className = 'ov zone look-' + (z.look || 'hover');
    setRect(d, z);
    if ((z.look || 'hover') === 'button') {
      d.textContent = z.label || '';
      if (z.color) d.style.background = z.color;
    }
    if (!editMode) {
      d.title = zoneTitle(z);
      d.addEventListener('click', function (ev) { ev.stopPropagation(); doAction(z); });
    } else {
      if (sel && sel.kind === 'zone' && sel.i === zi) { d.classList.add('sel'); addHandle(d); }
      attachEdit(d, 'zone', zi, z);
    }
    wrap.appendChild(d);
  });
  sizeButtons();
}
function addHandle(el) { var h = document.createElement('div'); h.className = 'hdl'; el.appendChild(h); }
function sizeButtons() {
  wrap.querySelectorAll('.look-button').forEach(function (el) {
    var hp = parseFloat(el.style.height) || 10;
    el.style.fontSize = Math.max(11, Math.round(wrap.clientHeight * hp / 100 * 0.38)) + 'px';
  });
}
function zoneTitle(z) {
  if (z.action === 'goto') return (FR ? 'Aller \u00e0 la diapo ' : 'Go to slide ') + ((z.slide || 0) + 1);
  if (z.action === 'url') return z.url || '';
  if (z.action === 'video') return FR ? 'Lire la vid\u00e9o' : 'Play video';
  if (z.action === 'back') return T.back;
  if (z.action === 'next') return FR ? 'Diapo suivante' : 'Next';
  if (z.action === 'prev') return FR ? 'Diapo pr\u00e9c\u00e9dente' : 'Previous';
  return '';
}
function doAction(z) {
  switch (z.action) {
    case 'goto': go(z.slide || 0); break;
    case 'next': go(linNext()); break;
    case 'prev': go(linPrev()); break;
    case 'back': goBack(); break;
    case 'url': if (z.url) window.open(z.url, '_blank'); break;
    case 'video': openLightbox(z.video); break;
  }
}

/* ---------- édition : sélection, déplacement, dessin ---------- */
function attachEdit(el, kind, i, obj) {
  el.addEventListener('pointerdown', function (e) {
    if (!editMode || drawMode) return;
    e.preventDefault();
    e.stopPropagation();
    select(kind, i);
    var p = relPct(e);
    drag = { mode: e.target.classList.contains('hdl') ? 'resize' : 'move',
             o: obj, el: wrap.querySelector('.ov.sel'),
             x0: p.x, y0: p.y, ox: obj.x, oy: obj.y, ow: obj.w, oh: obj.h };
    try { wrap.setPointerCapture(e.pointerId); } catch (err) {}
  });
}
function select(kind, i) { sel = { kind: kind, i: i }; renderOverlays(); renderProps(); }
function deselect() { sel = null; renderOverlays(); renderProps(); }

wrap.addEventListener('pointerdown', function (e) {
  if (!editMode) return;
  if (drawMode) {
    e.preventDefault();
    var p = relPct(e);
    var z = { x: p.x, y: p.y, w: 0, h: 0, action: 'goto',
              slide: Math.min(cur + 1, SLIDES.length - 1),
              look: 'button', label: 'Bouton', color: '#5b8cff' };
    SLIDES[cur].zones.push(z);
    sel = { kind: 'zone', i: SLIDES[cur].zones.length - 1 };
    renderOverlays();
    drag = { mode: 'draw', o: z, el: wrap.querySelector('.ov.sel'), x0: p.x, y0: p.y };
    try { wrap.setPointerCapture(e.pointerId); } catch (err) {}
  } else if (e.target === slideEl || e.target === wrap) deselect();
});
wrap.addEventListener('pointermove', function (e) {
  if (!drag) return;
  var p = relPct(e), o = drag.o;
  if (drag.mode === 'draw') {
    o.x = Math.min(p.x, drag.x0); o.y = Math.min(p.y, drag.y0);
    o.w = Math.abs(p.x - drag.x0); o.h = Math.abs(p.y - drag.y0);
  } else if (drag.mode === 'move') {
    o.x = clamp(drag.ox + (p.x - drag.x0), 0, 100 - o.w);
    o.y = clamp(drag.oy + (p.y - drag.y0), 0, 100 - o.h);
  } else {
    o.w = clamp(drag.ow + (p.x - drag.x0), 2, 100 - o.x);
    o.h = clamp(drag.oh + (p.y - drag.y0), 2, 100 - o.y);
  }
  if (drag.el) setRect(drag.el, o);
  markDirty();
});
function endDrag() {
  if (!drag) return;
  if (drag.mode === 'draw') {
    if (drag.o.w < 1.5 || drag.o.h < 1.5) { SLIDES[cur].zones.pop(); sel = null; renderOverlays(); }
    setDraw(false);
    renderProps();
  }
  sizeButtons();
  drag = null;
}
wrap.addEventListener('pointerup', endDrag);
wrap.addEventListener('pointercancel', endDrag);

function delSel() {
  var s = SLIDES[cur];
  if (!sel) return;
  if (sel.kind === 'zone') s.zones.splice(sel.i, 1);
  else { s.videos.splice(sel.i, 1); gcMedia(); }
  sel = null;
  markDirty();
  renderOverlays();
  renderProps();
}
function gcMedia() {
  var used = {};
  SLIDES.forEach(function (s) {
    s.videos.forEach(function (v) { if (v.media) used[v.media] = 1; });
    s.zones.forEach(function (z) { if (z.video && z.video.media) used[z.video.media] = 1; });
  });
  Object.keys(ASSETS.media).forEach(function (id) { if (!used[id]) delete ASSETS.media[id]; });
}

/* ---------- panneau propriétés ---------- */
function opt(v, label, curv) {
  return '<option value="' + escA(v) + '"' +
    (String(curv) === String(v) ? ' selected' : '') + '>' + esc(label) + '</option>';
}
function renderProps() {
  if (!editMode) { props.classList.add('hidden'); return; }
  props.classList.remove('hidden');
  var s = SLIDES[cur], z = null, v = null;
  var h = '<h3>Diapo ' + (cur + 1) + ' / ' + SLIDES.length + '</h3>' +
    '<label class="ck"><input type="checkbox" id="pHid"' + (s.hidden ? ' checked' : '') +
    '><span>Diapo cach\u00e9e \u2014 hors navigation, accessible uniquement via un bouton</span></label>';
  if (sel && sel.kind === 'zone' && s.zones[sel.i]) {
    z = s.zones[sel.i];
    h += '<hr><h3>Zone cliquable</h3>' +
      '<label>Action<select id="pAct">' +
      opt('goto', 'Aller \u00e0 une diapo', z.action) +
      opt('next', 'Diapo suivante', z.action) +
      opt('prev', 'Diapo pr\u00e9c\u00e9dente', z.action) +
      opt('back', 'Retour (diapo pr\u00e9c\u00e9demment vue)', z.action) +
      opt('url', 'Ouvrir un lien', z.action) +
      opt('video', 'Lire une vid\u00e9o (plein \u00e9cran)', z.action) +
      '</select></label>';
    if (z.action === 'goto') {
      h += '<label>Diapo cible<select id="pTgt">';
      SLIDES.forEach(function (t, i) {
        h += opt(i, 'Diapo ' + (i + 1) + (t.hidden ? ' (cach\u00e9e)' : ''), z.slide);
      });
      h += '</select></label>';
    }
    if (z.action === 'url')
      h += '<label>URL<input type="text" id="pUrl" value="' + escA(z.url || '') + '" placeholder="https://\u2026"></label>';
    if (z.action === 'video') {
      var vv = z.video || {};
      h += '<label>Source<select id="pVSrc">' + opt('yt', 'Lien YouTube', vv.media ? '' : 'yt');
      Object.keys(ASSETS.media).forEach(function (id) {
        h += opt(id, 'Vid\u00e9o embarqu\u00e9e ' + id, vv.media || '');
      });
      h += '</select></label>';
      if (!vv.media)
        h += '<label>Lien YouTube<input type="text" id="pVUrl" value="' + escA(vv.url || '') +
          '" placeholder="https://youtube.com/watch?v=\u2026"></label>';
    }
    h += '<label>Apparence<select id="pLook">' +
      opt('hover', 'Invisible (halo au survol)', z.look || 'hover') +
      opt('outline', 'Contour visible', z.look || 'hover') +
      opt('button', 'Bouton', z.look || 'hover') +
      '</select></label>';
    if ((z.look || 'hover') === 'button')
      h += '<label>Texte du bouton<input type="text" id="pLbl" value="' + escA(z.label || '') + '"></label>' +
        '<label>Couleur <input type="color" id="pCol" value="' + (z.color || '#5b8cff') + '"></label>';
    h += '<button class="danger" id="pDel">\ud83d\uddd1 Supprimer la zone</button>';
  } else if (sel && sel.kind === 'video' && s.videos[sel.i]) {
    v = s.videos[sel.i];
    h += '<hr><h3>Vid\u00e9o incrust\u00e9e</h3>';
    if (v.url) h += '<p class="muted">YouTube :<br>' + esc(v.url) + '</p>';
    else h += '<label class="ck"><input type="checkbox" id="pCtl"' + (v.controls !== false ? ' checked' : '') +
      '><span>Contr\u00f4les de lecture</span></label>' +
      '<label class="ck"><input type="checkbox" id="pAuto"' + (v.autoplay ? ' checked' : '') +
      '><span>Lecture auto (sans le son)</span></label>' +
      '<label class="ck"><input type="checkbox" id="pLoop"' + (v.loop ? ' checked' : '') +
      '><span>En boucle</span></label>';
    h += '<button class="danger" id="pDel">\ud83d\uddd1 Supprimer la vid\u00e9o</button>';
  } else {
    h += '<hr><p class="muted">' +
      'Clique sur une zone ou une vid\u00e9o de la diapo pour la modifier : ' +
      'd\u00e9place-la \u00e0 la souris, redimensionne avec la poign\u00e9e orange, <b>Suppr</b> pour l\u2019effacer.<br><br>' +
      '<b>\u2795 Zone</b> : dessine un bouton ou une zone cliquable.<br>' +
      '<b>\ud83c\udfac Vid\u00e9o</b> : incruste une vid\u00e9o sur la diapo.<br>' +
      '<b>\ud83d\udc41</b> sur une vignette : cache/montre la diapo.<br><br>' +
      '\ud83d\udcbe Enregistrer t\u00e9l\u00e9charge une copie \u00e0 jour de ce fichier HTML \u2014 ' +
      'remplace l\u2019ancien fichier par celle-ci.</p>';
  }
  props.innerHTML = h;
  var w = function (id, ev, fn) { var el = $(id); if (el) el.addEventListener(ev, fn); };
  w('pHid', 'change', function (e) { s.hidden = e.target.checked; markDirty(); structural(); });
  if (z) {
    w('pAct', 'change', function (e) {
      z.action = e.target.value;
      if (z.action === 'goto' && typeof z.slide !== 'number') z.slide = Math.min(cur + 1, SLIDES.length - 1);
      if (z.action === 'video' && !z.video) z.video = { url: '' };
      markDirty(); renderProps(); renderOverlays();
    });
    w('pTgt', 'change', function (e) { z.slide = parseInt(e.target.value, 10) || 0; markDirty(); renderOverlays(); });
    w('pUrl', 'input', function (e) { z.url = e.target.value.trim(); markDirty(); });
    w('pVSrc', 'change', function (e) {
      var val = e.target.value;
      z.video = val === 'yt' ? { url: '' } : { media: val };
      markDirty(); renderProps();
    });
    w('pVUrl', 'change', function (e) { z.video = { url: ytEmbed(e.target.value.trim()) }; markDirty(); renderProps(); });
    w('pLook', 'change', function (e) { z.look = e.target.value; markDirty(); renderProps(); renderOverlays(); });
    w('pLbl', 'input', function (e) { z.label = e.target.value; markDirty(); renderOverlays(); });
    w('pCol', 'input', function (e) { z.color = e.target.value; markDirty(); renderOverlays(); });
    w('pDel', 'click', delSel);
  } else if (v) {
    w('pCtl', 'change', function (e) { v.controls = e.target.checked; markDirty(); renderOverlays(); });
    w('pAuto', 'change', function (e) { v.autoplay = e.target.checked; if (v.autoplay) v.muted = true; markDirty(); renderOverlays(); });
    w('pLoop', 'change', function (e) { v.loop = e.target.checked; markDirty(); renderOverlays(); });
    w('pDel', 'click', delSel);
  }
}
function structural() { buildThumbs(); refresh(); }

/* ---------- vignettes ---------- */
function buildThumbs() {
  thumbs.innerHTML = '';
  thumbItems = [];
  SLIDES.forEach(function (s, i) {
    if (!editMode && s.hidden) return;
    var d = document.createElement('div');
    d.className = 'th' + (s.hidden ? ' th-hidden' : '');
    var im = document.createElement('img');
    im.src = IMG(s.img);
    im.draggable = false;
    d.appendChild(im);
    var n = document.createElement('span');
    n.className = 'tnum';
    n.textContent = i + 1;
    d.appendChild(n);
    if (editMode) {
      var b = document.createElement('button');
      b.className = 'teye';
      b.textContent = s.hidden ? '\ud83d\udeab' : '\ud83d\udc41';
      b.title = s.hidden ? 'Diapo cach\u00e9e \u2014 clic pour la remettre dans le fil'
        : 'Clic pour cacher cette diapo (accessible uniquement via un bouton)';
      b.addEventListener('click', function (ev) {
        ev.stopPropagation();
        s.hidden = !s.hidden;
        markDirty();
        structural();
      });
      d.appendChild(b);
    }
    d.addEventListener('click', function () { go(i); });
    thumbs.appendChild(d);
    thumbItems.push({ el: d, i: i });
  });
}
function syncThumbs() {
  thumbItems.forEach(function (t) {
    t.el.classList.toggle('current', t.i === cur);
    if (t.i === cur) t.el.scrollIntoView({ block: 'nearest' });
  });
}

/* ---------- ajout de vidéos ---------- */
function setDraw(on) {
  drawMode = on;
  $('btnAddZone').classList.toggle('active', on);
  document.body.classList.toggle('drawing', on);
}
function newMediaId() { var n = 1; while (ASSETS.media['m' + n]) n++; return 'm' + n; }
var vmenu = $('vmenu');
function hideVMenu() { vmenu.classList.add('hidden'); }
$('btnAddVideo').addEventListener('click', function (e) {
  e.stopPropagation();
  var r = e.currentTarget.getBoundingClientRect();
  vmenu.style.top = (r.bottom + 6) + 'px';
  vmenu.style.left = Math.max(8, r.right - 280) + 'px';
  vmenu.classList.toggle('hidden');
});
document.addEventListener('click', function (e) {
  if (!vmenu.classList.contains('hidden') && !vmenu.contains(e.target)) hideVMenu();
});
$('vmFile').addEventListener('click', function () { hideVMenu(); $('filePick').click(); });
$('filePick').addEventListener('change', function (e) {
  var f = e.target.files[0];
  e.target.value = '';
  if (!f) return;
  if (f.size > 60 * 1024 * 1024 &&
      !confirm('Fichier de ' + Math.round(f.size / 1e6) + ' Mo : le HTML final sera tr\u00e8s lourd. Continuer ?')) return;
  var r = new FileReader();
  r.onload = function () {
    var m = /^data:([^;]+);base64,([\s\S]*)$/.exec(r.result);
    if (!m) { alert('Lecture du fichier impossible.'); return; }
    var id = newMediaId();
    ASSETS.media[id] = { mime: m[1] || f.type || 'video/mp4', data: m[2] };
    SLIDES[cur].videos.push({ media: id, x: 25, y: 25, w: 50, h: 50, controls: true });
    markDirty();
    select('video', SLIDES[cur].videos.length - 1);
  };
  r.readAsDataURL(f);
});
$('vmYt').addEventListener('click', function () {
  hideVMenu();
  var u = prompt('Colle le lien YouTube :');
  if (!u) return;
  SLIDES[cur].videos.push({ url: ytEmbed(u.trim()), x: 25, y: 25, w: 50, h: 50 });
  markDirty();
  select('video', SLIDES[cur].videos.length - 1);
});
$('btnAddZone').addEventListener('click', function () { setDraw(!drawMode); });

/* ---------- lecteur vidéo plein écran ---------- */
function openLightbox(v) {
  if (!v) return;
  var lb = $('lb'), el = null;
  lb.innerHTML = '';
  if (v.url) {
    el = document.createElement('iframe');
    el.src = v.url + (v.url.indexOf('?') >= 0 ? '&' : '?') + 'autoplay=1';
    el.allow = 'autoplay; fullscreen; encrypted-media; picture-in-picture';
    el.allowFullscreen = true;
  } else if (v.media && ASSETS.media[v.media]) {
    el = document.createElement('video');
    el.src = MEDIA(v.media);
    el.controls = true;
    el.autoplay = true;
  }
  if (!el) return;
  lb.appendChild(el);
  $('lightbox').classList.remove('hidden');
  $('lbClose').classList.remove('hidden');
}
function closeLightbox() {
  $('lb').innerHTML = '';
  $('lightbox').classList.add('hidden');
  $('lbClose').classList.add('hidden');
}
$('lbClose').addEventListener('click', closeLightbox);
$('lightbox').addEventListener('click', function (e) { if (e.target.id === 'lightbox') closeLightbox(); });

/* ---------- enregistrement : le fichier se reconstruit lui-même ---------- */
function safeName(s) { return (s || 'presentation').replace(/[\\\/:*?"<>|]/g, '_'); }
function serialize(locked) {
  var cfg = JSON.parse(JSON.stringify(CFG));
  cfg.meta.locked = !!locked;
  var j = function (o) { return JSON.stringify(o).replace(/<\//g, '<\\/'); };
  return '<!DOCTYPE html>\n<html lang="' + META.lang + '">\n<head>\n<meta charset="utf-8">\n' +
    '<meta name="viewport" content="width=device-width, initial-scale=1">\n' +
    '<title>' + esc(META.title) + '</title>\n</head>\n<body>\n' +
    '<script type="application/json" id="cfg">' + j(cfg) + '<\/script>\n' +
    '<script type="application/json" id="assets">' + j(ASSETS) + '<\/script>\n' +
    '<script id="app-src">' + $('app-src').textContent + '<\/script>\n' +
    '</body>\n</html>';
}
function download(name, txt) {
  var b = new Blob([txt], { type: 'text/html' });
  var a = document.createElement('a');
  a.href = URL.createObjectURL(b);
  a.download = name;
  document.body.appendChild(a);
  a.click();
  a.remove();
  setTimeout(function () { URL.revokeObjectURL(a.href); }, 10000);
}
$('btnSave').addEventListener('click', function () {
  download(safeName(META.title) + '.html', serialize(false));
  clearDirty();
});
$('btnLock').addEventListener('click', function () {
  if (!confirm('Exporter une version VERROUILL\u00c9E (sans mode \u00e9dition), pr\u00eate \u00e0 diffuser ?\n' +
    'Garde aussi ta version normale : c\u2019est elle qui reste modifiable.')) return;
  download(safeName(META.title) + ' \u2014 final.html', serialize(true));
});
window.addEventListener('beforeunload', function (e) {
  if (dirty) { e.preventDefault(); e.returnValue = ''; }
});

/* ---------- mode édition ---------- */
function setEdit(on) {
  if (META.locked) return;
  editMode = on;
  sel = null;
  setDraw(false);
  document.body.classList.toggle('editing', on);
  $('tools').classList.toggle('hidden', !on);
  var be = $('btnEdit');
  if (be) be.classList.toggle('active', on);
  if (on) { thumbs.classList.remove('hidden'); $('btnThumbs').classList.add('active'); }
  buildThumbs();
  refresh();
  renderProps();
}
var btnEdit = $('btnEdit');
if (btnEdit) btnEdit.addEventListener('click', function () { setEdit(!editMode); });
titleEl.addEventListener('dblclick', function () {
  if (!editMode) return;
  var t = prompt('Titre de la pr\u00e9sentation :', META.title);
  if (t && t.trim()) { META.title = t.trim(); document.title = META.title; markDirty(); }
});

/* ---------- boutons d'entête, clavier, tactile ---------- */
if (SLIDES.some(function (s) { return s.notes; })) btnNotes.classList.remove('hidden');
btnNotes.addEventListener('click', function () {
  notesEl.classList.toggle('hidden');
  btnNotes.classList.toggle('active', !notesEl.classList.contains('hidden'));
});
$('btnThumbs').addEventListener('click', function () {
  thumbs.classList.toggle('hidden');
  $('btnThumbs').classList.toggle('active', !thumbs.classList.contains('hidden'));
});
$('btnFS').addEventListener('click', function () {
  document.fullscreenElement ? document.exitFullscreen()
    : document.documentElement.requestFullscreen();
});
$('prev').addEventListener('click', function () { if (!editMode) go(linPrev()); });
$('next').addEventListener('click', function () { if (!editMode) go(linNext()); });
backBtn.addEventListener('click', goBack);

document.addEventListener('keydown', function (e) {
  if (e.target.closest && e.target.closest('input,textarea,select')) return;
  var k = e.key;
  if (k === 'Escape') {
    if (!$('lightbox').classList.contains('hidden')) { closeLightbox(); return; }
    if (drawMode) { setDraw(false); return; }
    if (editMode && sel) { deselect(); return; }
    return;
  }
  if (editMode && (k === 'Delete' || k === 'Backspace')) {
    if (sel) { e.preventDefault(); delSel(); }
    return;
  }
  if (k === 'ArrowRight' || k === 'PageDown' || k === ' ') {
    e.preventDefault();
    if (editMode) go(Math.min(cur + 1, SLIDES.length - 1));
    else if (!SLIDES[cur].hidden) go(linNext());
  } else if (k === 'ArrowLeft' || k === 'PageUp') {
    if (editMode) go(Math.max(cur - 1, 0));
    else if (SLIDES[cur].hidden) goBack();
    else go(linPrev());
  } else if (k === 'Home') go(editMode ? 0 : firstVisible());
  else if (k === 'End') go(editMode ? SLIDES.length - 1 : lastVisible());
  else if (k.toLowerCase() === 'f') $('btnFS').click();
  else if (k.toLowerCase() === 't') $('btnThumbs').click();
  else if (k.toLowerCase() === 'n' && !btnNotes.classList.contains('hidden')) btnNotes.click();
  else if (k.toLowerCase() === 'e') setEdit(!editMode);
});

var tx = null;
document.addEventListener('touchstart', function (e) { tx = e.touches[0].clientX; }, { passive: true });
document.addEventListener('touchend', function (e) {
  if (tx === null || editMode || !$('lightbox').classList.contains('hidden')) { tx = null; return; }
  var dx = e.changedTouches[0].clientX - tx;
  tx = null;
  if (Math.abs(dx) < 50) return;
  if (SLIDES[cur].hidden) { if (dx > 0) goBack(); return; }
  go(dx < 0 ? linNext() : linPrev());
});

window.addEventListener('resize', sizeButtons);
slideEl.addEventListener('load', sizeButtons);
window.addEventListener('hashchange', function () {
  var h = parseInt(location.hash.slice(1), 10) - 1;
  if (!isNaN(h) && h !== cur) go(clamp(h, 0, SLIDES.length - 1), { noHist: true });
});

/* ---------- démarrage ---------- */
setTimeout(function () {
  var h = $('hint');
  if (h) { h.style.opacity = 0; setTimeout(function () { h.remove(); }, 600); }
}, 5000);
buildThumbs();
$('btnThumbs').classList.add('active');
var start = parseInt(location.hash.slice(1), 10) - 1;
go(isNaN(start) ? firstVisible() : clamp(start, 0, SLIDES.length - 1),
   { instant: true, noHist: true });
})();
"""


def write_html(out_path, title, slides_src, notes, zones, videos, media,
               lang="fr", embed=True):
    cfg = {"meta": {"title": title, "lang": lang, "embed": embed,
                    "locked": False, "app": APP_VERSION},
           "slides": [{"img": i, "notes": notes[i], "hidden": False,
                       "zones": zones[i], "videos": videos[i]}
                      for i in range(len(slides_src))]}
    assets = {"images": slides_src, "media": media}

    def j(o):
        # '</' échappé pour ne jamais fermer la balise <script> qui porte le JSON
        return json.dumps(o, ensure_ascii=False,
                          separators=(",", ":")).replace("</", "<\\/")

    doc = ("<!DOCTYPE html>\n"
           '<html lang="%s">\n<head>\n<meta charset="utf-8">\n'
           '<meta name="viewport" content="width=device-width, initial-scale=1">\n'
           "<title>%s</title>\n</head>\n<body>\n"
           '<script type="application/json" id="cfg">%s</script>\n'
           '<script type="application/json" id="assets">%s</script>\n'
           '<script id="app-src">%s</script>\n'
           "</body>\n</html>"
           % (lang, html.escape(title), j(cfg), j(assets), APP_JS))
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(doc)


# ============================== INTERFACE ==============================

LO_URL = "https://fr.libreoffice.org/download/telecharger-libreoffice/"
PPTX_EXTS = (".pptx", ".ppt", ".potx")


def ensure_deps_gui():
    """Premier lancement en mode script : installe les bibliothèques
    manquantes via pip, avec une petite fenêtre d'attente. Permet de se
    passer d'un installeur .bat (bloqué par Smart App Control)."""
    if getattr(sys, "frozen", False):
        return
    need = []
    try:
        import pptx  # noqa: F401
    except ImportError:
        need.append("python-pptx")
    if get_pymupdf() is None:
        need.append("pymupdf")
    try:
        import tkinterdnd2  # noqa: F401
    except ImportError:
        need.append("tkinterdnd2")
    if not need:
        return
    import threading
    import tkinter as tk
    w = tk.Tk()
    w.title("Première installation")
    w.geometry("460x130")
    w.configure(bg="#111318")
    tk.Label(w, text="Installation des composants (première fois, ~1 min)…\n"
             + ", ".join(need), bg="#111318", fg="#e8eaf0",
             font=("Segoe UI", 10), justify="center").pack(expand=True)

    def work():
        kw = {}
        if sys.platform == "win32":
            kw["creationflags"] = subprocess.CREATE_NO_WINDOW
        subprocess.run([sys.executable, "-m", "pip", "install", "--user",
                        "--quiet"] + need, capture_output=True, **kw)
        w.after(0, w.destroy)

    threading.Thread(target=work, daemon=True).start()
    w.mainloop()
    import importlib
    importlib.invalidate_caches()


def launch_gui(preload=None):
    import threading, webbrowser
    import tkinter as tk
    from tkinter import filedialog

    # glisser-déposer si tkinterdnd2 est disponible (toujours le cas dans l'exe)
    try:
        from tkinterdnd2 import DND_FILES, TkinterDnD
        root = TkinterDnD.Tk()
        has_dnd = True
    except Exception:
        root = tk.Tk()
        has_dnd = False

    BG, PANEL, FG, MUTED, ACCENT = "#111318", "#1b1e26", "#e8eaf0", "#8b90a0", "#5b8cff"

    root.title(f"PPTX → HTML interactif — v{APP_VERSION}")
    root.geometry("560x560")
    root.configure(bg=BG)
    root.minsize(480, 500)

    state = {"queue": [], "busy": False}

    tk.Label(root, text="PPTX → HTML interactif", bg=BG, fg=FG,
             font=("Segoe UI", 16, "bold")).pack(pady=(18, 2))
    tk.Label(root, text="Google Slides : Fichier → Télécharger → Microsoft PowerPoint (.pptx)\n"
             "Le HTML généré contient un mode édition (✏️) : boutons, diapos cachées, vidéos.",
             bg=BG, fg=MUTED, font=("Segoe UI", 10), justify="center").pack(pady=(0, 12))

    # --- zone de dépôt ---
    dropbox = tk.Frame(root, bg=PANEL, highlightbackground="#2a2e3a",
                       highlightthickness=1, cursor="hand2")
    dropbox.pack(fill="x", padx=24, pady=4)
    drop_lbl = tk.Label(
        dropbox,
        text="⬇\nDépose ton fichier .pptx ici" if has_dnd
        else "Clique pour choisir un fichier .pptx",
        bg=PANEL, fg=FG, font=("Segoe UI", 12, "bold"), justify="center",
        cursor="hand2")
    drop_lbl.pack(pady=(18, 2))
    file_lbl = tk.Label(
        dropbox,
        text="(ou clique pour parcourir — plusieurs fichiers acceptés)"
        if has_dnd else "(plusieurs fichiers acceptés)",
        bg=PANEL, fg=MUTED, font=("Segoe UI", 9), cursor="hand2")
    file_lbl.pack(pady=(0, 16))

    # --- options ---
    opts = tk.Frame(root, bg=BG)
    opts.pack(fill="x", padx=24, pady=(12, 4))
    assets_var = tk.BooleanVar(value=False)
    tk.Checkbutton(opts, text="Images dans un dossier séparé (le HTML ne "
                   "fonctionne pas s'il est copié seul — anti-diffusion)",
                   variable=assets_var, bg=BG, fg=FG, selectcolor=PANEL,
                   activebackground=BG, activeforeground=FG,
                   font=("Segoe UI", 9), wraplength=480,
                   justify="left").pack(anchor="w")
    qrow = tk.Frame(opts, bg=BG); qrow.pack(anchor="w", pady=(8, 0))
    tk.Label(qrow, text="Qualité :", bg=BG, fg=FG,
             font=("Segoe UI", 9)).pack(side="left")
    dpi_var = tk.StringVar(value="150")
    for label, val in [("Légère (100)", "100"), ("Standard (150)", "150"),
                       ("Haute (200)", "200")]:
        tk.Radiobutton(qrow, text=label, value=val, variable=dpi_var, bg=BG,
                       fg=FG, selectcolor=PANEL, activebackground=BG,
                       activeforeground=FG, font=("Segoe UI", 9)
                       ).pack(side="left", padx=6)

    # --- bouton + log ---
    go_btn = tk.Button(root, text="Convertir", state="disabled", bg="#2e3342",
                       fg="white", relief="flat", cursor="hand2",
                       font=("Segoe UI", 12, "bold"), padx=30, pady=8)
    go_btn.pack(pady=14)
    logbox = tk.Text(root, height=6, bg=PANEL, fg=MUTED, relief="flat",
                     font=("Consolas", 9), state="disabled", wrap="word")
    logbox.pack(fill="both", expand=True, padx=24, pady=(0, 18))

    def log(msg):
        def _():
            logbox.config(state="normal")
            logbox.insert("end", msg + "\n")
            logbox.see("end")
            logbox.config(state="disabled")
        root.after(0, _)

    def open_folder(path):
        d = os.path.dirname(path)
        try:
            if sys.platform == "win32":
                os.startfile(d)
            elif sys.platform == "darwin":
                subprocess.Popen(["open", d])
            else:
                subprocess.Popen(["xdg-open", d])
        except Exception:
            pass

    def start_convert():
        if state["busy"] or not state["queue"]:
            return
        state["busy"] = True
        go_btn.config(state="disabled", text="Conversion…", bg="#2e3342")

        def work():
            done, last = 0, None
            for f in list(state["queue"]):
                out = os.path.splitext(f)[0] + ".html"
                try:
                    log(f"▶ {os.path.basename(f)}")
                    last = convert(f, out, dpi=int(dpi_var.get()),
                                   embed=not assets_var.get(), log=log)
                    done += 1
                except ConvError as e:
                    log("ERREUR : " + str(e))
                except Exception as e:
                    log("ERREUR inattendue : " + repr(e))
            if done and last:
                open_folder(last)
            state["busy"] = False
            root.after(0, lambda: go_btn.config(
                state="normal", text="Convertir", bg=ACCENT))
        threading.Thread(target=work, daemon=True).start()

    go_btn.config(command=start_convert)

    def set_files(paths, auto=True):
        if state["busy"]:
            log("Conversion en cours — dépose ton fichier quand c'est terminé.")
            return
        valid = [os.path.abspath(p) for p in paths
                 if os.path.isfile(p) and p.lower().endswith(PPTX_EXTS)]
        skipped = len(paths) - len(valid)
        if skipped:
            log(f"{skipped} fichier(s) ignoré(s) — seuls "
                f"{', '.join(PPTX_EXTS)} sont acceptés.")
        if not valid:
            return
        state["queue"] = valid
        file_lbl.config(fg=FG, text=os.path.basename(valid[0])
                        if len(valid) == 1 else f"{len(valid)} fichiers sélectionnés")
        go_btn.config(state="normal", bg=ACCENT)
        if auto:
            start_convert()

    def choose(_event=None):
        p = filedialog.askopenfilenames(
            title="Choisir une ou plusieurs présentations",
            filetypes=[("PowerPoint", "*.pptx *.ppt *.potx"), ("Tous", "*.*")])
        if p:
            set_files(list(p))

    for w in (dropbox, drop_lbl, file_lbl):
        w.bind("<Button-1>", choose)

    if has_dnd:
        def on_drop(e):
            set_files(list(root.tk.splitlist(e.data)))
        for w in (root, dropbox, drop_lbl, file_lbl):
            w.drop_target_register(DND_FILES)
            w.dnd_bind("<<Drop>>", on_drop)

    def check_deps():
        try:
            find_soffice()
            log("LibreOffice : OK")
        except ConvError:
            log("LibreOffice manquant — il sert à rendre les slides "
                "(installation unique, tout reste hors ligne).")
            tk.Button(root, text="Installer LibreOffice",
                      command=lambda: webbrowser.open(LO_URL),
                      bg="#c0392b", fg="white", relief="flat", cursor="hand2",
                      font=("Segoe UI", 10, "bold"), padx=14, pady=5
                      ).pack(pady=(0, 10))
        if get_pymupdf() is not None:
            log("Rendu d'images : PyMuPDF OK")
        else:
            try:
                find_pdftoppm()
                log("Rendu d'images : Poppler OK")
            except ConvError as e:
                log("ERREUR : " + str(e))
        log("Prêt — dépose un .pptx, la conversion démarre toute seule."
            if has_dnd else "Prêt — clique dans la zone pour choisir un .pptx.")

    check_deps()
    if preload:
        root.after(300, lambda: set_files([preload]))
    root.mainloop()


if __name__ == "__main__":
    argv = sys.argv[1:]
    if argv and len(argv) == 1 and os.path.isfile(argv[0]) \
            and argv[0].lower().endswith(PPTX_EXTS):
        # fichier déposé sur l'icône (exe ou .pyw) → GUI préchargée
        ensure_deps_gui()
        launch_gui(preload=os.path.abspath(argv[0]))
    elif argv:
        ap = argparse.ArgumentParser()
        ap.add_argument("pptx")
        ap.add_argument("-o", "--output", default=None)
        ap.add_argument("--dpi", type=int, default=150)
        ap.add_argument("--lang", default="fr", choices=["fr", "en"])
        ap.add_argument("--assets", action="store_true")
        a = ap.parse_args()
        out = a.output or os.path.splitext(a.pptx)[0] + ".html"
        try:
            convert(a.pptx, out, dpi=a.dpi, lang=a.lang, embed=not a.assets)
        except ConvError as e:
            sys.exit(str(e))
    else:
        ensure_deps_gui()
        launch_gui()
